import os
import json
import base64
from openai import AsyncOpenAI
from typing import List, Dict

client = AsyncOpenAI(api_key=os.getenv("OPENAI_API_KEY"))

IMAGE_EXTENSIONS = {'.jpg', '.jpeg', '.png', '.gif', '.webp'}
MIME_TYPES = {
    '.jpg': 'image/jpeg', '.jpeg': 'image/jpeg',
    '.png': 'image/png', '.gif': 'image/gif', '.webp': 'image/webp',
}


def extract_text_from_file(file_path: str, file_ext: str) -> str:
    """
    Extract plain text from a document file.
    Supports: xlsx/xls, docx/doc, pdf, pptx/ppt
    Returns a plain-text string with the content.
    """
    if file_ext in ('.xlsx', '.xls'):
        return read_excel_file(file_path)
    elif file_ext in ('.docx', '.doc'):
        return read_word_file(file_path)
    elif file_ext == '.pdf':
        return read_pdf_file(file_path)
    elif file_ext in ('.pptx', '.ppt'):
        return read_pptx_file(file_path)
    else:
        raise ValueError(f"Unsupported document extension: {file_ext}")


def file_to_base64_data_uri(file_path: str, file_ext: str) -> str:
    """Convert an image file to a base64 data URI for GPT-4o vision."""
    mime = MIME_TYPES.get(file_ext, 'image/jpeg')
    with open(file_path, 'rb') as f:
        encoded = base64.b64encode(f.read()).decode('utf-8')
    return f"data:{mime};base64,{encoded}"


async def extract_content_from_uploaded_files(
    file_paths: List[tuple]  # list of (file_path, file_ext)
) -> Dict:
    """
    Process a list of uploaded files and return:
      - extracted_text: combined text from all document files
      - images: list of base64 data URIs for image files
    """
    text_parts = []
    images = []

    for file_path, file_ext in file_paths:
        try:
            if file_ext in IMAGE_EXTENSIONS:
                images.append(file_to_base64_data_uri(file_path, file_ext))
            else:
                text = extract_text_from_file(file_path, file_ext)
                if text.strip():
                    text_parts.append(f"--- File: {os.path.basename(file_path)} ---\n{text}")
        except Exception as e:
            # Don't fail the whole request if one file can't be parsed
            text_parts.append(f"--- File: {os.path.basename(file_path)} (could not parse: {e}) ---")

    return {
        'extracted_text': '\n\n'.join(text_parts),
        'images': images,
    }


async def parse_course_file(file_path: str, file_ext: str, teacher_uid: str):
    """
    Parse uploaded course file using GPT-4 to extract structured data
    """

    # Read file content based on type
    if file_ext in ['.xlsx', '.xls']:
        file_content = read_excel_file(file_path)
    elif file_ext in ['.docx', '.doc']:
        file_content = read_word_file(file_path)
    else:
        raise ValueError(f"Unsupported file extension: {file_ext}")

    # Use GPT-4 to parse and structure the content
    prompt = f"""
You are a curriculum parser. Your job is to extract EVERY piece of content from the teacher's file and map it faithfully into a structured course format — losing nothing.

FILE CONTENT:
{file_content}

OUTPUT FORMAT (JSON only, no markdown):
{{
  "course_name": "string",
  "class": "string (e.g., Class 1, Grade 3 — infer from content if not stated)",
  "subject": "string (e.g., Math, Science, Health — infer if not stated)",
  "topic": "string (main theme of the course)",
  "sections": [
    {{
      "id": "section-1",
      "title": "string (e.g., the day/unit/chapter name)",
      "description": "string (overview of what this section covers)",
      "subsections": [
        {{
          "id": "sub-1-1",
          "title": "string (logical grouping within the section, e.g., 'Main Topics', 'Activities', 'Key Concepts')",
          "description": "string (what this group covers)",
          "topicBoxes": [
            {{
              "id": "topic-1-1-1",
              "title": "string (ONE specific topic, concept, organ, activity, or skill from the file)",
              "description": "string (2-3 sentences: what is covered, any activity instructions or notes from the file go here)",
              "duration_minutes": 20,
              "pla_pillars": [],
              "learning_objectives": ["string — infer 1-2 measurable objectives per topic"],
              "content_keywords": ["keyword1", "keyword2"],
              "video_resources": [],
              "worksheets": [],
              "activities": ["string — copy any activity or game description from the file here"]
            }}
          ]
        }}
      ]
    }}
  ]
}}

CRITICAL RULES — read carefully:

1. CAPTURE EVERYTHING. Every listed topic, organ, concept, activity, game, or note in the file must appear as a topicBox somewhere. Do not summarize or skip items.

2. STRUCTURE: Map the file's natural hierarchy:
   - Top-level groups (days, units, chapters) → sections
   - Sub-groupings (Main Topics, Activities, Videos) → subsections
   - Individual items (each organ, each activity, each concept) → topicBoxes

3. TOPIC BOXES are the atomic unit. Each topicBox should cover exactly ONE item:
   - "Brain", "Heart", "Liver" → three separate topicBoxes, not one
   - "Hand-washing demo", "Hygiene matching game" → two separate topicBoxes
   - "Sing Head, Shoulders, Knees & Toes" → one topicBox

4. ACTIVITIES: If the file lists activities, games, demos, or worksheets:
   - Each activity becomes its own topicBox under an "Activities" subsection
   - Copy the activity description into the topicBox's "description" field
   - Also put it in the "activities" array
   - If it's a song, game, craft, or demo — make that clear in the description

5. LINKS/URLs: If any links or URLs appear in the file, put them in:
   - "video_resources" if they are video links
   - "worksheets" if they are worksheet/document links
   - "activities" if they are activity/game links

6. SUBSECTION GROUPING: Group topicBoxes logically:
   - Organs/body parts → one subsection (e.g., "Key Organs & Body Parts")
   - Activities & games → another subsection (e.g., "Activities & Games")
   - Notes or additional content → another subsection if needed

7. IDs: Use section-1, section-2, sub-1-1, sub-1-2, topic-1-1-1, topic-1-1-2, etc.

8. INFER when needed:
   - Grade/class: infer from content complexity
   - Subject: infer from content domain
   - Learning objectives: write 1-2 measurable objectives per topicBox (use action verbs: identify, describe, demonstrate, explain, list)
   - Duration: 10-15 min per topicBox is typical

9. Return ONLY valid JSON. No markdown, no explanations, no code fences.
"""

    response = await client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": "You are a curriculum structure parser. Output only valid JSON. Capture every single item from the source file as a topicBox — never skip or summarize content."},
            {"role": "user", "content": prompt}
        ],
        temperature=0.2,
        max_tokens=16000,
        response_format={"type": "json_object"}
    )

    parsed_text = response.choices[0].message.content.strip()
    course_data = json.loads(parsed_text)

    return course_data


def read_excel_file(file_path: str) -> str:
    """Read Excel file and convert to text, preserving merged cell structure"""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True)
        content_parts = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            content_parts.append(f"SHEET: {sheet_name}\n")

            # Build a cell value map that fills merged cell ranges with the merged value
            merged_values = {}
            for merged_range in ws.merged_cells.ranges:
                top_left = ws.cell(merged_range.min_row, merged_range.min_col).value
                for row in merged_range.rows:
                    for cell in row:
                        merged_values[(cell.row, cell.column)] = top_left

            rows_text = []
            for row in ws.iter_rows():
                row_values = []
                for cell in row:
                    val = merged_values.get((cell.row, cell.column), cell.value)
                    row_values.append(str(val) if val is not None else "")
                # Skip completely empty rows
                if any(v.strip() for v in row_values):
                    rows_text.append(" | ".join(row_values))

            content_parts.append("\n".join(rows_text))
            content_parts.append("\n\n")

        return "\n".join(content_parts)

    except ImportError:
        # Fallback to pandas if openpyxl not available directly
        import pandas as pd
        try:
            excel_file = pd.ExcelFile(file_path)
            content_parts = []
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name, header=None)
                content_parts.append(f"SHEET: {sheet_name}\n")
                # Forward-fill to handle merged cells
                df = df.ffill(axis=0)
                for _, row in df.iterrows():
                    row_vals = [str(v) for v in row if str(v) not in ('nan', 'None', '')]
                    if row_vals:
                        content_parts.append(" | ".join(row_vals))
                content_parts.append("\n\n")
            return "\n".join(content_parts)
        except Exception as e:
            raise ValueError(f"Failed to read Excel file: {str(e)}")
    except Exception as e:
        raise ValueError(f"Failed to read Excel file: {str(e)}")


def read_word_file(file_path: str) -> str:
    """Read Word document and convert to text"""
    from docx import Document

    try:
        doc = Document(file_path)
        content_parts = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                content_parts.append(paragraph.text)
        return "\n".join(content_parts)
    except Exception as e:
        raise ValueError(f"Failed to read Word file: {str(e)}")


def read_pdf_file(file_path: str) -> str:
    """Extract text from a PDF using pypdf"""
    try:
        from pypdf import PdfReader
        reader = PdfReader(file_path)
        pages = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ''
            if text.strip():
                pages.append(f"[Page {i+1}]\n{text}")
        return "\n\n".join(pages)
    except ImportError:
        raise ValueError("pypdf is not installed. Add 'pypdf' to requirements.txt.")
    except Exception as e:
        raise ValueError(f"Failed to read PDF: {str(e)}")


def read_pptx_file(file_path: str) -> str:
    """Extract text from a PowerPoint file using python-pptx"""
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(slides)
    except ImportError:
        raise ValueError("python-pptx is not installed. Add 'python-pptx' to requirements.txt.")
    except Exception as e:
        raise ValueError(f"Failed to read PowerPoint: {str(e)}")